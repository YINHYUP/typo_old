import os
from pptx import Presentation

def extract_text_from_pptx(file_path):
    text = []
    presentation = Presentation(file_path)
    
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    
    return "\n".join(text)

def main(pptx_folder, output_file):
    all_text = []
    
    for filename in os.listdir(pptx_folder):
        if filename.endswith('.pptx'):
            file_path = os.path.join(pptx_folder, filename)
            text = extract_text_from_pptx(file_path)
            all_text.append(f"--- {filename} ---\n{text}\n")
    
    with open(output_file, 'w', encoding='utf-8') as f:
        f.writelines(all_text)

if __name__ == "__main__":
    pptx_folder = 'D:\#.Secure Work Folder\python'  # 파워포인트 파일이 있는 폴더 경로
    output_file = 'output.txt'  # 저장할 텍스트 파일 이름
    main(pptx_folder, output_file)

